import streamlit as st
import io
import pandas as pd
import numpy as np
import fitz  # PyMuPDF
from pptx import Presentation
from sklearn.cluster import KMeans
from sklearn.preprocessing import StandardScaler
from sklearn.decomposition import PCA
import plotly.express as px

st.title('🧠 Customer Segmentation App with File Support')

uploaded_file = st.file_uploader(
    "📁 Upload CSV, Excel, PDF, or PPTX file", 
    type=["csv", "xlsx", "pdf", "pptx"]
)

def extract_pdf_text(file):
    text = ""
    doc = fitz.open(stream=file.read(), filetype="pdf")
    for page in doc:
        text += page.get_text()
    return text

def extract_ppt_text(file):
    text = ""
    prs = Presentation(file)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

if uploaded_file is not None:
    file_name = uploaded_file.name.lower()

    if file_name.endswith('.csv'):
        decoded = pd.read_csv(uploaded_file)
        file_type = 'table'
    elif file_name.endswith('.xlsx'):
        decoded = pd.read_excel(uploaded_file)
        file_type = 'table'
    elif file_name.endswith('.pdf'):
        st.subheader("📄 Extracted Text from PDF")
        st.text(extract_pdf_text(uploaded_file))
        file_type = 'text'
    elif file_name.endswith('.pptx'):
        st.subheader("📽️ Extracted Text from PPTX")
        st.text(extract_ppt_text(uploaded_file))
        file_type = 'text'
    else:
        st.error("❌ Unsupported file format.")
        st.stop()

    if file_type == 'table':
        st.subheader("📊 Data Preview")
        st.dataframe(decoded.head())

        # Select and clean numeric features
        features = decoded.select_dtypes(include=[np.number])
        features = features.replace([np.inf, -np.inf], np.nan)

        # Keep only rows with valid numeric data
        valid_idx = features.dropna().index
        features = features.loc[valid_idx].reset_index(drop=True)
        decoded = decoded.loc[valid_idx].reset_index(drop=True)

        if features.empty:
            st.warning("⚠️ No usable numeric data after cleaning.")
            st.stop()

        st.write(f"✅ Clustering on {features.shape[0]} valid rows.")

        # Cluster selection
        num_clusters = st.slider("🎯 Select number of clusters", 2, 10, 3)

        def perform_clustering(data, n_clusters):
            scaler = StandardScaler()
            scaled_data = scaler.fit_transform(data)
            kmeans = KMeans(n_clusters=n_clusters, random_state=42)
            return kmeans.fit_predict(scaled_data)

        clusters = perform_clustering(features, num_clusters)

        # Add clusters to original data
        decoded['Cluster'] = clusters

        # PCA for 2D visualization
        pca = PCA(n_components=2)
        reduced_data = pca.fit_transform(features)

        fig = px.scatter(
            x=reduced_data[:, 0],
            y=reduced_data[:, 1],
            color=clusters.astype(str),
            title='🌀 Customer Segmentation (PCA Visualized)',
            labels={'x': 'PCA 1', 'y': 'PCA 2', 'color': 'Cluster'},
        )
        fig.update_layout(margin=dict(l=20, r=20, t=40, b=20))

        st.plotly_chart(fig)

        st.subheader("📋 Clustered Data Sample")
        st.dataframe(decoded.head(10))
    else:
        st.info("ℹ️ Only CSV and Excel files are used for clustering. Text files shown for reference.")
